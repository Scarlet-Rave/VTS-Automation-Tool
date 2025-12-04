import os
import json
import re
import uuid
import tempfile
from time import sleep
from io import BytesIO
from copy import deepcopy
from math import ceil

import random
import streamlit as st
import httpx
from dotenv import load_dotenv
from twelvelabs import TwelveLabs
from twelvelabs.indexes import IndexesCreateRequestModelsItem
from docx import Document
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
import pyodbc
import pandas as pd
from sklearn.linear_model import LinearRegression
from sklearn.ensemble import RandomForestClassifier
from sklearn.preprocessing import LabelEncoder
import plotly.express as px

# ====================
# 🔹 Environment Setup
# ====================
load_dotenv()
API_KEY = os.getenv("TWELVE_API_KEY")
if not API_KEY:
    st.error("❌ Missing TWELVE_API_KEY in your .env file.")
    st.stop()

st.set_page_config(page_title="Video Converter Automation", layout="centered")
st.title("📄 VTS Automation Tool")
with st.expander("⚠️ Important Notes"):
    st.markdown(
        "- Ensure your video has at least **360p resolution**.\n"
        "- Longer videos take longer to process.\n"
    )

uploaded_file = st.file_uploader("Upload Source Video", type=["mp4","mov","avi","mkv"])
if uploaded_file and st.session_state.get("uploaded_name") != uploaded_file.name:
    st.session_state.doc_ready = False
    st.session_state.doc_path = None
    st.session_state.pptx_ready = False
    st.session_state.pptx_stream = None
    st.session_state.show_pptx_download = True
    st.session_state.steps_with_difficulty = []
    st.session_state.uploaded_name = uploaded_file.name
    st.session_state.displayed_ui = False

st.subheader("📌 Process Assembly Template Preview")
folder = os.path.join(os.path.dirname(__file__), "Templates")
col1, col2 = st.columns(2)
with col1:
    img1 = os.path.join(folder, "preview_concept_1.png")
    if os.path.exists(img1): st.image(img1, use_container_width=True)
with col2:
    img2 = os.path.join(folder, "preview_concept_2.png")
    if os.path.exists(img2): st.image(img2, use_container_width=True)

template_option = st.selectbox("Choose a PowerPoint template:", ["Concept 1","Concept 2"])
template_files = {"Concept 1": os.path.join(folder,"Concept_1.pptx"),
                  "Concept 2": os.path.join(folder,"Concept_2.pptx")}
template_path = template_files[template_option]
if not os.path.exists(template_path):
    st.error(f"Template file not found at {template_path}")
    st.stop()

output_name_input = st.text_input("Output filename (without .pptx)", value="procedure_document")

http_client = httpx.Client(verify=False, timeout=httpx.Timeout(300.0))
client = TwelveLabs(api_key=API_KEY, httpx_client=http_client)

# ====================
# 🔹 Session State
# ====================
if "progress" not in st.session_state:
    st.session_state.progress = 0
if "status" not in st.session_state:
    st.session_state.status = "⏳ Waiting to start..."
if "doc_ready" not in st.session_state:
    st.session_state.doc_ready = False
if "doc_path" not in st.session_state:
    st.session_state.doc_path = None
if "pptx_ready" not in st.session_state:
    st.session_state.pptx_ready = False
if "pptx_stream" not in st.session_state:
    st.session_state.pptx_stream = None
if "show_pptx_download" not in st.session_state:
    st.session_state.show_pptx_download = True
if "temp_dir" not in st.session_state:
    st.session_state.temp_dir = None
if "steps_with_difficulty" not in st.session_state:
    st.session_state.steps_with_difficulty = []
if "displayed_ui" not in st.session_state:
    st.session_state.displayed_ui = False

progress_bar = None
progress_text = None

# ====================
# 🔹 Progress Helper
# ====================
def update_progress(value,message=""):
    if progress_bar and progress_text:
        st.session_state.progress = value
        st.session_state.status = message
        progress_bar.progress(value)
        progress_text.text(f"{value}% - {message}")
        sleep(0.05)

# ====================
# 🔹 TwelveLabs Indexing
# ====================
def create_index_if_needed(name="video-proc-index"):
    update_progress(10,"Creating or fetching index...")
    try:
        index = client.indexes.create(
            index_name=name,
            models=[IndexesCreateRequestModelsItem(model_name="pegasus1.2",model_options=["visual","audio"])]
        )
        update_progress(20,"Index created")
        return index.id
    except Exception:
        idxs = client.indexes.list()
        for idx in idxs:
            if getattr(idx,"index_name","")==name:
                update_progress(20,"Using existing index")
                return idx.id
        raise RuntimeError("Failed to create or fetch index")

def upload_and_index(video_path,index_id):
    update_progress(25,"Uploading and indexing video...")
    with open(video_path,"rb") as f:
        task = client.tasks.create(index_id=index_id, video_file=f, request_options={"timeout":300})
    task_id = task.id
    update_progress(40,"Indexing video...")
    task = client.tasks.wait_for_done(task_id=task_id)
    if task.status!="ready":
        raise RuntimeError("Indexing failed: "+str(task.status))
    # 🔹 Correct video_id extraction
    if hasattr(task, "video_id") and task.video_id:
        return task.video_id
    elif hasattr(task, "output") and isinstance(task.output, dict) and "video_id" in task.output:
        return task.output["video_id"]
    elif hasattr(task, "output") and isinstance(task.output, list):
        for o in task.output:
            if isinstance(o, dict) and "video_id" in o:
                return o["video_id"]
    raise RuntimeError("No valid video_id returned from TwelveLabs task")

# ====================
# 🔹 Procedural Steps Extraction
# ====================
def ask_for_steps_bilingual(video_id):
    update_progress(70,"Extracting procedural steps...")
    prompt="""
Extract the procedural steps from this video.
Each step should include both English and Malay versions.
DO NOT translate the word "Step".
Each step must START with "Step X" in both the english and malay strings.
Return strict JSON only (no extra text). Format exactly as:

[
  {
    "step": 1,
    "english": "Step 1 Insert the motherboard into the computer case, aligning the back panel connectors with the I.O. shield.",
    "malay": "Step 1 Masukkan papan induk ke dalam casing komputer, sejajarkan penyambung panel belakang dengan I.O. shield.",
    "start_time": 0.0,
    "end_time": 10.5
  },
  ...
]
Do not include tools or commentary.
"""
    res = client.analyze(video_id=video_id, prompt=prompt)
    text = getattr(res,"data",None) or getattr(res,"output",None) or str(res)
    m = re.search(r'(\[.*\])',text,flags=re.DOTALL)
    json_text = m.group(1) if m else text
    try:
        steps = json.loads(json_text)
    except Exception as e:
        raise RuntimeError("Failed to parse JSON from model output. Raw output:\n"+text) from e
    for s in steps:
        stepnum = s.get("step","")
        step_str = str(stepnum)
        eng = s.get("english","")
        mal = s.get("malay","")
        if not eng.startswith(f"Step {step_str}"):
            s["english"]=f"Step {step_str} "+eng.lstrip()
        mal_clean = re.sub(r'^\s*Step\s*\d+\s*','',mal,flags=re.IGNORECASE)
        s["malay"]=f"Step {step_str} "+mal_clean.lstrip()
    update_progress(85,"Extracted steps")
    return steps

# ====================
# 🔹 Word Export
# ====================
def write_docx_bilingual(steps,out_path):
    update_progress(90,"🖋 Writing Word document (English + Malay)...")
    doc = Document()
    doc.add_heading("Procedural Steps (English + Malay)",level=1)
    for s in steps:
        stepnum = s.get("step","")
        step_str = str(stepnum)
        eng = s.get("english","").strip()
        mal = s.get("malay","").strip()
        if not eng.startswith(f"Step {step_str}"):
            eng = f"Step {step_str} "+eng
        mal = re.sub(r'^\s*Step\s*\d+\s*','',mal,flags=re.IGNORECASE)
        mal = f"Step {step_str} "+mal
        start = s.get("start_time","")
        end = s.get("end_time","")
        doc.add_paragraph(eng)
        doc.add_paragraph(mal)
        if start!="" or end!="":
            doc.add_paragraph(f"⏱ Timestamp: {start} — {end}")
        doc.add_paragraph("")
    doc.save(out_path)
    update_progress(100,"Document ready")
    return out_path

def parse_docx_steps(docx_stream):
    doc = Document(docx_stream)
    steps_by_num = {}
    seen_counts = {}
    def ensure_step(n):
        if n not in steps_by_num:
            steps_by_num[n]={"step":n,"english":"","malay":"","timestamp":"","tools":""}
            seen_counts[n]=0
        return steps_by_num[n]
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if text.startswith("⏱") or text.startswith("🧰"):
            if steps_by_num:
                last_step_num = max(steps_by_num.keys())
                if text.startswith("⏱"):
                    steps_by_num[last_step_num]["timestamp"]=text
                else:
                    steps_by_num[last_step_num]["tools"]=text
            continue
        m = re.match(r"Step\s*(\d+)\b(.*)",text,flags=re.IGNORECASE)
        if m:
            n = int(m.group(1))
            rest = m.group(2).strip()
            step = ensure_step(n)
            count = seen_counts.get(n,0)
            if count==0:
                step["english"]=f"Step {n} {rest}" if rest else f"Step {n}"
                seen_counts[n]=1
            elif count==1:
                step["malay"]=f"Step {n} {rest}" if rest else f"Step {n}"
                seen_counts[n]=2
            else:
                seen_counts[n]+=1
            continue
        if steps_by_num:
            last_step_num=max(steps_by_num.keys())
            step=steps_by_num[last_step_num]
            if not step["english"]:
                step["english"]=text
            elif not step["malay"]:
                step["malay"]=text
            else:
                step["malay"]+=" "+text
    steps=[steps_by_num[k] for k in sorted(steps_by_num.keys())]
    return steps

# ====================
# 🔹 PowerPoint Functions
# ====================
STEP_RE=re.compile(r"SECTION:Step(\d+)(EN|BM|TS)")

def get_existing_step_numbers(prs):
    found=set()
    for slide in prs.slides:
        for shape in slide.shapes:
            name=getattr(shape,"name","")
            m=STEP_RE.search(name)
            if m:
                found.add(int(m.group(1)))
    return sorted(found)

def slide_min_step(slide):
    nums=[]
    for shape in slide.shapes:
        name=getattr(shape,"name","")
        m=STEP_RE.search(name)
        if m:
            nums.append(int(m.group(1)))
    return min(nums) if nums else None

def slide_contains_step(slide,step_num):
    for shape in slide.shapes:
        name=getattr(shape,"name","")
        if name==f"SECTION:Step{step_num}EN":
            return True
    return False

def duplicate_slide(prs,source_slide):
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    sp_tree=new_slide.shapes._spTree
    for shape in source_slide.shapes:
        el=shape.element
        new_el=deepcopy(el)
        sp_tree.insert_element_before(new_el,'p:extLst')
    return new_slide

def shift_slide_placeholders(slide,shift):
    for shape in slide.shapes:
        try:
            name=getattr(shape,"name","")
            if not name:
                continue
            m=STEP_RE.search(name)
            if m:
                old_n=int(m.group(1))
                suffix=m.group(2)
                new_n=old_n+shift
                shape.name=f"SECTION:Step{new_n}{suffix}"
        except Exception:
            pass

def fill_pptx(template_stream,steps):
    prs=Presentation(template_stream)
    existing=get_existing_step_numbers(prs)
    if not existing:
        raise ValueError("Template contains no SECTION:Step placeholders")
    max_template_step=max(existing)
    max_needed_step=max(step["step"] for step in steps) if steps else 0
    if max_needed_step>max_template_step:
        extra_needed=max_needed_step-max_template_step
        slides_to_add=ceil(extra_needed/4)
        base_slide=None
        for s in prs.slides:
            if slide_contains_step(s,max_template_step):
                base_slide=s
                break
        if base_slide is None:
            for s in prs.slides:
                if slide_min_step(s) is not None:
                    base_slide=s
                    break
        if base_slide is None:
            raise ValueError("Could not find suitable base slide.")
        for i in range(1,slides_to_add+1):
            new_slide=duplicate_slide(prs,base_slide)
            shift_slide_placeholders(new_slide,i*4)
    replacements={}
    for step in steps:
        n=step["step"]
        replacements[f"SECTION:Step{n}EN"]=step["english"]
        replacements[f"SECTION:Step{n}BM"]=step["malay"]
        replacements[f"SECTION:Step{n}TS"]="\n".join(filter(None,[step.get("timestamp",""),step.get("tools","")]))
    step_counter=0
    for slide in prs.slides:
        step_placeholders=[]
        for shape in slide.shapes:
            name=getattr(shape,"name","")
            if re.match(r"SECTION:Step\d+$",name):
                step_placeholders.append(shape)
        step_placeholders.sort(key=lambda s:s.name)
        for shape in step_placeholders:
            step_counter+=1
            tf=shape.text_frame
            tf.clear()
            p=tf.paragraphs[0]
            run=p.add_run()
            run.text=f"Step {step_counter}"
            font=run.font
            font.name="Times New Roman"
            font.size=Pt(18)
            font.bold=True
            font.color.rgb=RGBColor(0x00,0x70,0xC0)
        for shape in slide.shapes:
            name=getattr(shape,"name","")
            if hasattr(shape,"has_text_frame") and shape.has_text_frame and STEP_RE.match(name):
                text=replacements.get(name,"")
                shape.text_frame.clear()
                if text:
                    shape.text_frame.text=text
                is_malay=name.endswith("BM")
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.name="Times New Roman"
                        run.font.size=Pt(14)
                        run.font.color.rgb=RGBColor(0,122,204) if is_malay else RGBColor(0,0,0)
                try:
                    shape.fill.background()
                    shape.line.fill.background()
                except:
                    pass
                shape.name=""
            if hasattr(shape,"has_table") and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        key=cell.text.strip()
                        if key in replacements:
                            text=replacements[key]
                            cell.text=text
                            try:
                                cell.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
                            except:
                                pass
                            is_malay=key.endswith("BM")
                            for p in cell.text_frame.paragraphs:
                                try:
                                    p.alignment=PP_ALIGN.CENTER
                                except:
                                    pass
                                for run in p.runs:
                                    run.font.name="Times New Roman"
                                    run.font.size=Pt(14)
                                    run.font.color.rgb=RGBColor(0,122,204) if is_malay else RGBColor(0,0,0)
    if "Concept_2" in getattr(template_stream,"name","") or "Concept 2" in str(template_stream):
        if prs.slides:
            last_slide=prs.slides[-1]
            for shape in list(last_slide.shapes):
                if getattr(shape,"name","")=="SECTION:DownArrow":
                    try:
                        shape.element.getparent().remove(shape.element)
                    except:
                        pass
    output_stream=BytesIO()
    prs.save(output_stream)
    output_stream.seek(0)
    return output_stream

# ====================
# 🔹 Risk Analysis with TwelveLabs
# ====================
def generate_risk_and_precaution(step_text, video_id):
    """
    Use TwelveLabs to generate risk level, extra precaution, and estimated time for a step.
    """
    prompt = f"""
Analyze this procedural step: "{step_text}".
Return a strict JSON like:
{{
    "risk_level": "Low|Medium|High",
    "precaution": "string",
    "estimated_time": number (in minutes)
}}
"""
    res = client.analyze(video_id=video_id, prompt=prompt)
    text = getattr(res,"data",None) or getattr(res,"output",None) or str(res)
    m = re.search(r'(\{.*\})', text, flags=re.DOTALL)
    json_text = m.group(1) if m else text
    try:
        data = json.loads(json_text)
    except Exception:
        # fallback
        data = {"risk_level":"Low","precaution":"","estimated_time":1}
    risk = data.get("risk_level","Low")
    precaution = data.get("precaution","")
    est_time = float(data.get("estimated_time",1))
    return risk,precaution,est_time

# ====================
# 🔹 Access DB
# ====================
def save_steps_to_access(steps, video_name=None):
    db_path=os.path.join(os.path.dirname(__file__),"Database","database.accdb")
    if not os.path.exists(db_path):
        raise FileNotFoundError(f"Access DB not found at {db_path}")
    conn_str=r'DRIVER={Microsoft Access Driver (*.mdb, *.accdb)};DBQ='+db_path
    conn=pyodbc.connect(conn_str)
    cursor=conn.cursor()
    table_name="ProcedureSteps"
    create_sql=f"""
    CREATE TABLE [{table_name}] (
        [ID] AUTOINCREMENT PRIMARY KEY,
        [VideoName] TEXT(255),
        [StepNum] INT,
        [RiskLevel] TEXT(50),
        [ExtraPrecaution] TEXT(255),
        [EstimatedTime] DOUBLE
    )
    """
    try:
        cursor.execute(create_sql)
    except:
        pass
    existing_steps = {}
    if video_name:
        cursor.execute(f"SELECT StepNum,RiskLevel,ExtraPrecaution,EstimatedTime FROM [{table_name}] WHERE VideoName=?",video_name)
        rows = cursor.fetchall()
        for r in rows:
            existing_steps[r.StepNum]={"risk_level":r.RiskLevel,"precaution":r.ExtraPrecaution,"estimated_time":r.EstimatedTime}
    for s in steps:
        n = s["step"]
        if n in existing_steps:
            s["risk_level"]=existing_steps[n]["risk_level"]
            s["precaution"]=existing_steps[n]["precaution"]
            s["estimated_time"]=existing_steps[n]["estimated_time"]
        else:
            # 🔹 Use actual TwelveLabs data
            risk,precaution,est_time = generate_risk_and_precaution(s.get("english",""), st.session_state.video_id)
            s["risk_level"]=risk
            s["precaution"]=precaution
            s["estimated_time"]=est_time
        cursor.execute(f"""
            INSERT INTO [{table_name}] (VideoName,[StepNum],[RiskLevel],[ExtraPrecaution],[EstimatedTime])
            VALUES (?,?,?,?,?)
        """, video_name, s["step"], s["risk_level"], s["precaution"], s["estimated_time"])
    conn.commit()
    conn.close()
    return table_name

# ====================
# 🔹 Predict Difficulty
# ====================
def predict_difficulty_regression(steps):
    df=pd.DataFrame(steps)
    if df.empty:
        return steps
    df["risk_num"]=df["risk_level"].map({"Low":1,"Medium":2,"High":3})
    X=df[["risk_num","estimated_time"]].fillna(0)
    y=(df.index+1).astype(float)
    model=LinearRegression()
    model.fit(X,y)
    preds=model.predict(X)
    df["predicted_difficulty"]=preds
    for idx,s in enumerate(steps):
        s["predicted_difficulty"]=round(float(df.loc[idx,"predicted_difficulty"]),2)
    return steps

# ====================
# 🔹 Display UI
# ====================
def display_steps_interactive(steps):
    st.subheader("📌 Procedural Steps Overview")
    risk_counts={"Low":0,"Medium":0,"High":0}
    for s in steps:
        risk = s.get("risk_level","Low")
        risk_counts[risk] = risk_counts.get(risk,0)+1
    fig_pie=px.pie(
        names=list(risk_counts.keys()),
        values=list(risk_counts.values()),
        title="Risk Level Distribution",
        color=list(risk_counts.keys()),
        color_discrete_map={"Low":"green","Medium":"orange","High":"red"}
    )
    st.plotly_chart(fig_pie,use_container_width=True,key=f"risk_chart_{uuid.uuid4()}")
    df_duration=pd.DataFrame([{"Step":f"Step {s['step']}","Duration":s["estimated_time"],"Risk":s.get("risk_level","Low")} for s in steps])
    per_bar_height=40
    total_chart_height=max(400,per_bar_height*len(steps))
    fig_duration=px.bar(
        df_duration,
        y="Step",
        x="Duration",
        color="Risk",
        color_discrete_map={"Low":"green","Medium":"orange","High":"red"},
        title="Step Duration Overview",
        orientation="h",
        height=total_chart_height
    )
    max_duration=max([s["estimated_time"] for s in steps]+[5])
    fig_duration.update_layout(
        xaxis=dict(range=[0,max_duration]),
        yaxis=dict(autorange="reversed"),
        bargap=0.2,
        uniformtext_minsize=10,
        uniformtext_mode="hide"
    )
    st.markdown('<div style="overflow-y:auto; max-height:400px;">',unsafe_allow_html=True)
    st.plotly_chart(fig_duration,use_container_width=True,key=f"duration_chart_{uuid.uuid4()}")
    st.markdown('</div>',unsafe_allow_html=True)
    search_query=st.text_input("🔍 Search for Step (press Enter to search)")
    for s in steps:
        step_text=f"Step {s.get('step')}"
        if search_query.strip() and search_query.strip().lower() not in step_text.lower():
            continue
        difficulty=s.get("predicted_difficulty",0)
        risk=s.get("risk_level","Medium")
        if risk=="Low":
            card_color="#d4f4dd"
        elif risk=="Medium":
            card_color="#fff1b8"
        else:
            card_color="#ffc6c4"
        with st.expander(f"{step_text}",expanded=bool(search_query.strip())):
            st.markdown(
                f"""
                <div style="background-color:{card_color}; border-radius:8px; padding:8px; margin-bottom:5px;">
                    <table style="width:100%; border-collapse: collapse;">
                        <tr>
                            <td style="padding:4px; font-weight:bold; width:150px;">English</td>
                            <td style="padding:4px;">{s.get('english','')}</td>
                        </tr>
                        <tr>
                            <td style="padding:4px; font-weight:bold;">Malay</td>
                            <td style="padding:4px;">{s.get('malay','')}</td>
                        </tr>
                        <tr>
                            <td style="padding:4px; font-weight:bold;">Risk Level</td>
                            <td style="padding:4px;">{s.get('risk_level','')}</td>
                        </tr>
                        <tr>
                            <td style="padding:4px; font-weight:bold;">Extra Precaution</td>
                            <td style="padding:4px;">{s.get('precaution','')}</td>
                        </tr>
                        <tr>
                            <td style="padding:4px; font-weight:bold;">Estimated Time</td>
                            <td style="padding:4px;">{s.get('estimated_time',0)} min</td>
                        </tr>
                        <tr>
                            <td style="padding:4px; font-weight:bold;">Predicted Difficulty</td>
                            <td style="padding:4px; color:{'red' if difficulty>3 else 'green'}">{difficulty}</td>
                        </tr>
                    </table>
                </div>
                """,unsafe_allow_html=True
            )

# ====================
# 🔹 Download Helper
# ====================
def _hide_download_button():
    st.session_state.show_pptx_download=False

# ====================
# 🔹 Main Processing Logic
# ====================
if uploaded_file and output_name_input:
    if st.button("Start Processing"):
        progress_bar=st.progress(0)
        progress_text=st.empty()
        temp_dir=tempfile.mkdtemp()
        st.session_state.temp_dir=temp_dir
        video_path=os.path.join(temp_dir,uploaded_file.name)
        with open(video_path,"wb") as f:
            f.write(uploaded_file.getbuffer())
        try:
            index_id=create_index_if_needed()
            video_id=upload_and_index(video_path,index_id)
            st.session_state.video_id = video_id  # 🔹 Save video_id in session
            steps=ask_for_steps_bilingual(video_id)
            docx_path=os.path.join(temp_dir,f"{output_name_input}.docx")
            write_docx_bilingual(steps,docx_path)
            st.session_state.doc_ready=True
            st.session_state.doc_path=docx_path
            with open(docx_path,"rb") as f:
                steps_parsed=parse_docx_steps(f)
            table_name=save_steps_to_access(steps_parsed, video_name=uploaded_file.name)
            steps_with_difficulty=predict_difficulty_regression(steps_parsed)
            st.session_state.steps_with_difficulty=steps_with_difficulty
            with open(template_path,"rb") as f:
                pptx_stream=fill_pptx(f,steps_with_difficulty)
            st.session_state.pptx_ready=True
            st.session_state.pptx_stream=pptx_stream
            st.session_state.show_pptx_download=True
            st.session_state.displayed_ui=True
            st.success("Conversion completed!")
        except Exception as e:
            st.error(f"❌ Error: {str(e)}")

if st.session_state.get("displayed_ui"):
    display_steps_interactive(st.session_state.steps_with_difficulty)
    if st.session_state.get("pptx_stream") and st.session_state.get("show_pptx_download",True):
        st.download_button(label="📥 Download Converted File",
                           data=st.session_state.pptx_stream,
                           file_name=f"{output_name_input}.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                           on_click=_hide_download_button,
                           key=f"download_btn_{uuid.uuid4()}")
